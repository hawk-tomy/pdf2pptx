# Copyright (c) 2020 Kevin McGuinness <kevin.mcguinness@gmail.com>
import io
from pathlib import Path
from typing import TYPE_CHECKING

import fitz  # type: ignore
from pptx import Presentation
from pptx.util import Cm
from tqdm import trange

if TYPE_CHECKING:
    from fitz import Document  # type: ignore

__all__ = ['convert_pdf2pptx']


def convert_pdf2pptx(
    pdf_file: Path,
    output_file: str | None,
    resolution: int,
    start_page: int,
    page_count: int | None,
    quiet: bool = False,
    note_file: str | None = None,
):
    doc: Document = fitz.open(pdf_file)
    if not quiet:
        print(pdf_file, 'contains', doc.page_count, 'slides')  # type: ignore

    path = Path(pdf_file)
    if note_file is None:
        note_path = path.with_name(path.name.replace('.pdf', '_notes.txt'))
    else:
        note_path = Path(note_file)

    _notes: list[str] | None = None
    try:
        if note_path.exists():
            _notes = ['']
            for line in note_path.read_text().splitlines(keepends=True):
                if line.startswith('---'):
                    _notes.append('')
                    continue
                if line.startswith('//'):
                    continue
                _notes[-1] += line
            if not _notes[-1]:
                _notes.pop()
    except Exception:
        pass
    finally:
        notes: list[str] = _notes or []

    if page_count is None:
        page_count = doc.page_count  # type: ignore
    assert isinstance(page_count, int)

    # transformation matrix: slide to pixmap
    zoom = resolution / 72
    matrix = fitz.Matrix(zoom, zoom, 0)

    # create pptx presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # configure presentation aspect ratio
    page = doc.load_page(0)  # type: ignore
    aspect_ratio = page.rect.width / page.rect.height
    prs.slide_width = int(prs.slide_height * aspect_ratio)  # type: ignore

    # create page iterator
    if not quiet:
        page_iter = trange(start_page, start_page + page_count)
    else:
        page_iter = range(start_page, start_page + page_count)

    # iterate over slides
    for page_no in page_iter:
        page = doc.load_page(page_no)  # type: ignore

        # write slide as a pixmap
        pixmap = page.get_pixmap(matrix=matrix)  # type: ignore
        image_data = pixmap.tobytes(output='PNG')  # type: ignore
        image_file = io.BytesIO(image_data)  # type: ignore

        # add a slide
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Cm(0)
        slide.shapes.add_picture(image_file, left, top, height=prs.slide_height)

        note = notes[page_no] if page_no < len(notes) else ''
        slide.notes_slide.notes_text_frame.text += note  # type: ignore

    if output_file is None:
        output_file = path.with_suffix('.pptx')  # type: ignore

    # save presentation
    prs.save(output_file)  # type: ignore
